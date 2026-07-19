import os
import uuid
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from config import Config
from database.db_manager import get_db_connection

class DocumentProcessor:
    @staticmethod
    def extract_text(file_path: str, ext: str) -> str:
        text = ""
        if ext == ".pdf":
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
        elif ext == ".docx":
            doc = Document(file_path)
            for p in doc.paragraphs:
                text += p.text + "\n"
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        return text

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200):
        chunks = []
        words = text.split()
        stride = chunk_size - chunk_overlap
        for i in range(0, len(words), stride):
            chunk = " ".join(words[i:i + chunk_size])
            chunks.append(chunk)
            if i + chunk_size >= len(words):
                break
        return chunks