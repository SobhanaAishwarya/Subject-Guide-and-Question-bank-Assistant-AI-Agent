"""
Document ingestion: extract → clean → chunk.

Supports PDF, DOCX, PPTX and TXT. Page (or slide) numbers are preserved all the
way through to the chunk metadata so answers can cite them.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Callable, List, Optional, Sequence, Tuple

from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import settings
from models.schemas import Chunk
from utils.logger import get_logger
from utils.text_utils import clean_text

logger = get_logger(__name__)

# A page/slide unit: (page_number_or_None, text)
PageUnit = Tuple[Optional[int], str]


class UnsupportedFileType(ValueError):
    """Raised when an uploaded file extension is not supported."""


class DocumentProcessor:
    """Turns uploaded files into clean, embeddable :class:`Chunk` objects."""

    def __init__(
        self,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
    ) -> None:
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size or settings.chunk_size,
            chunk_overlap=chunk_overlap or settings.chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", "? ", "! ", "; ", ", ", " ", ""],
        )

    # ------------------------------------------------------------------ #
    # Extraction
    # ------------------------------------------------------------------ #
    @staticmethod
    def _extract_pdf(data: bytes) -> List[PageUnit]:
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages: List[PageUnit] = []
        for index, page in enumerate(reader.pages, start=1):
            try:
                pages.append((index, page.extract_text() or ""))
            except Exception as exc:  # noqa: BLE001 - a bad page must not kill the doc
                logger.warning("Failed to extract PDF page %s: %s", index, exc)
                pages.append((index, ""))
        return pages

    @staticmethod
    def _extract_docx(data: bytes) -> List[PageUnit]:
        import docx

        document = docx.Document(io.BytesIO(data))
        parts: List[str] = [p.text for p in document.paragraphs if p.text.strip()]

        # Tables often carry definitions, so flatten them into readable rows.
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        # DOCX has no reliable page concept, so the whole file is one unit.
        return [(None, "\n".join(parts))]

    @staticmethod
    def _extract_pptx(data: bytes) -> List[PageUnit]:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        slides: List[PageUnit] = []
        for index, slide in enumerate(presentation.slides, start=1):
            fragments: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        fragments.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            fragments.append(" | ".join(cells))
            slides.append((index, "\n".join(fragments)))
        return slides

    @staticmethod
    def _extract_txt(data: bytes) -> List[PageUnit]:
        for encoding in ("utf-8", "utf-16", "latin-1"):
            try:
                return [(None, data.decode(encoding))]
            except UnicodeDecodeError:
                continue
        return [(None, data.decode("utf-8", errors="ignore"))]

    def extract(self, filename: str, data: bytes) -> List[PageUnit]:
        """Extract raw per-page text from a file's bytes."""
        extension = Path(filename).suffix.lower().lstrip(".")
        extractors: dict[str, Callable[[bytes], List[PageUnit]]] = {
            "pdf": self._extract_pdf,
            "docx": self._extract_docx,
            "pptx": self._extract_pptx,
            "txt": self._extract_txt,
        }
        if extension not in extractors:
            raise UnsupportedFileType(
                f"'{extension}' is not supported. Use: "
                f"{', '.join(settings.supported_extensions)}."
            )
        logger.info("Extracting %s (%s)", filename, extension)
        return extractors[extension](data)

    # ------------------------------------------------------------------ #
    # Chunking
    # ------------------------------------------------------------------ #
    def chunk_pages(
        self,
        pages: Sequence[PageUnit],
        source: str,
        subject: str = "General",
        doc_id: str = "",
    ) -> List[Chunk]:
        """Clean each page and split it into overlapping chunks."""
        chunks: List[Chunk] = []
        counter = 0

        for page_number, raw_text in pages:
            cleaned = clean_text(raw_text)
            if len(cleaned) < 30:  # skip blank or near-blank pages
                continue
            for piece in self.splitter.split_text(cleaned):
                piece = piece.strip()
                if len(piece) < 30:
                    continue
                chunk = Chunk(
                    text=piece,
                    source=source,
                    page=page_number,
                    chunk_index=counter,
                    subject=subject,
                )
                if doc_id:
                    chunk.doc_id = doc_id
                chunks.append(chunk)
                counter += 1

        logger.info("Produced %s chunks from %s", len(chunks), source)
        return chunks

    def process(
        self,
        filename: str,
        data: bytes,
        subject: str = "General",
        doc_id: str = "",
    ) -> Tuple[List[Chunk], int]:
        """
        Full pipeline for one file.

        Returns ``(chunks, page_count)``.
        """
        pages = self.extract(filename, data)
        chunks = self.chunk_pages(pages, source=filename, subject=subject, doc_id=doc_id)
        return chunks, len(pages)
